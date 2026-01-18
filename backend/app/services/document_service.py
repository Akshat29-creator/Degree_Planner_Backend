"""
Document Processing Service for Revision Analyzer.

Extracts text from PDF and PPT files for AI analysis.
"""
import io
from typing import Tuple

# PyMuPDF (fitz module)
import fitz

# python-pptx
from pptx import Presentation


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from a PDF file."""
    text = ""
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        raise ValueError(f"Failed to extract text from PDF: {e}")
    return text.strip()


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from a PowerPoint file."""
    text = ""
    try:
        prs = Presentation(io.BytesIO(file_content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {e}")
    return text.strip()


def extract_text(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    Detect file type and extract text.
    Returns (extracted_text, file_type)
    """
    lower_filename = filename.lower()
    if lower_filename.endswith(".pdf"):
        return extract_text_from_pdf(file_content), "pdf"
    elif lower_filename.endswith(".pptx") or lower_filename.endswith(".ppt"):
        return extract_text_from_pptx(file_content), "pptx"
    else:
        raise ValueError(f"Unsupported file type: {filename}. Only PDF and PPTX are supported.")
